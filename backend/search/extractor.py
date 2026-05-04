"""
openclaw/backend/search/extractor.py

Extracts plain text from uploaded documents.
Supports: PDF, DOCX, PPTX, XLSX, TXT, MD, and plain fallback.

Usage:
    from backend.search.extractor import extract_text
    text = extract_text(path="/path/to/file.pdf", mime_type="application/pdf")
"""
import logging
from pathlib import Path

logger = logging.getLogger("openclaw.search.extractor")


def extract_text(path: str, mime_type: str = "") -> str:
    """
    Extract all readable text from a file.
    Returns empty string if extraction fails or format unsupported.
    Never raises — logs warning on failure.
    """
    p = Path(path)
    if not p.exists():
        logger.warning(f"[extractor] File not found: {path}")
        return ""

    suffix = p.suffix.lower()

    try:
        # ── PDF ───────────────────────────────────────────────────────────────
        if suffix == ".pdf" or "pdf" in mime_type:
            return _extract_pdf(p)

        # ── Word ──────────────────────────────────────────────────────────────
        if suffix in (".docx", ".doc") or "word" in mime_type or "officedocument.wordprocessingml" in mime_type:
            return _extract_docx(p)

        # ── PowerPoint ────────────────────────────────────────────────────────
        if suffix in (".pptx", ".ppt") or "presentationml" in mime_type:
            return _extract_pptx(p)

        # ── Excel ─────────────────────────────────────────────────────────────
        if suffix in (".xlsx", ".xls") or "spreadsheetml" in mime_type:
            return _extract_xlsx(p)

        # ── Plain text / Markdown ─────────────────────────────────────────────
        if suffix in (".txt", ".md", ".rst", ".csv", ".json", ".yaml", ".yml"):
            return p.read_text(encoding="utf-8", errors="replace")

        # ── Unknown — try reading as text ─────────────────────────────────────
        logger.info(f"[extractor] Unknown type {suffix}, attempting plain text read")
        return p.read_text(encoding="utf-8", errors="replace")

    except Exception as exc:
        logger.warning(f"[extractor] Failed to extract {path}: {exc}")
        return ""


# ── Format-specific extractors ────────────────────────────────────────────────

def _extract_pdf(path: Path) -> str:
    import fitz  # PyMuPDF
    doc  = fitz.open(str(path))
    pages = []
    for i, page in enumerate(doc):
        text = page.get_text("text")
        if text.strip():
            pages.append(f"[Page {i+1}]\n{text.strip()}")
    doc.close()
    return "\n\n".join(pages)


def _extract_docx(path: Path) -> str:
    from docx import Document
    doc    = Document(str(path))
    parts  = []

    # Body paragraphs
    for para in doc.paragraphs:
        if para.text.strip():
            # Include heading level as context
            style = para.style.name if para.style else ""
            prefix = f"[{style}] " if "Heading" in style else ""
            parts.append(f"{prefix}{para.text.strip()}")

    # Tables
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))

    return "\n\n".join(parts)


def _extract_pptx(path: Path) -> str:
    from pptx import Presentation
    prs   = Presentation(str(path))
    slides = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        if texts:
            slides.append(f"[Slide {i+1}]\n" + "\n".join(texts))
    return "\n\n".join(slides)


def _extract_xlsx(path: Path) -> str:
    import openpyxl
    wb    = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    sheets = []
    for sheet in wb.worksheets:
        rows = []
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                rows.append(" | ".join(cells))
        if rows:
            sheets.append(f"[Sheet: {sheet.title}]\n" + "\n".join(rows))
    wb.close()
    return "\n\n".join(sheets)
